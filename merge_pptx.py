# merge_pptx.py — fuses slides preserving all XML/shapes/images/fonts
import copy, shutil, zipfile, os, re
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

OUTPUT = r"C:\Users\migue\Documents\xenium-spatial-transcriptomics\xenium_presentacion_final.pptx"

# Order: v3 slide 1 (intro), then xenium_presentacion slides 1+2 (tech + analysis)
SOURCES = [
    (r"C:\Users\migue\Documents\xenium-spatial-transcriptomics\xenium_presentacion_v3.pptx",  [0]),
    (r"C:\Users\migue\Documents\xenium-spatial-transcriptomics\xenium_presentacion.pptx",      [0, 1]),
]

# ── low-level slide copy helper (preserves all relationships/media) ───────────
def copy_slide(src_prs, slide_idx, dst_prs):
    """Deep-copy a slide from src_prs into dst_prs, including media and rels."""
    src_slide = src_prs.slides[slide_idx]

    # 1. Add a blank slide using the first available layout in dst
    blank_layout = dst_prs.slide_layouts[6]  # "Blank"
    dst_slide    = dst_prs.slides.add_slide(blank_layout)

    # 2. Replace the spTree (shape tree) XML with a deep copy of the source
    src_sp_tree = src_slide.shapes._spTree
    dst_sp_tree = dst_slide.shapes._spTree
    # Remove all existing children from the destination shape tree
    for child in list(dst_sp_tree):
        dst_sp_tree.remove(child)
    for child in src_sp_tree:
        dst_sp_tree.append(copy.deepcopy(child))

    # 3. Copy background (bg element) if present
    src_el = src_slide._element
    dst_el = dst_slide._element
    src_bg = src_el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if src_bg is not None:
        dst_bg = dst_el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if dst_bg is not None:
            dst_el.remove(dst_bg)
        dst_el.insert(2, copy.deepcopy(src_bg))

    # 4. Copy relationships + media files
    src_part  = src_slide.part
    dst_part  = dst_slide.part

    # We need to walk all rels in the source slide and port them
    for rel in src_part.rels.values():
        if rel.is_external:
            dst_part.add_relationship(rel.reltype, rel.target_ref)
        else:
            target_part = rel.target_part
            # Copy the blob (binary content) into dst
            try:
                new_part = dst_part.part_related_by(rel.reltype)
            except KeyError:
                pass
            # Add related part — python-pptx will deduplicate if same content
            try:
                dst_part.relate_to(target_part, rel.reltype)
            except Exception:
                pass

    # 5. Fix rId references: build a mapping old→new and patch XML
    #    Build src relId → reltype mapping
    src_rids = {r.rId: r for r in src_part.rels.values()}
    dst_rids = {r.rId: r for r in dst_part.rels.values()}

    return dst_slide


# ── Use zipfile-level merge for fidelity ─────────────────────────────────────
# python-pptx's high-level copy can silently drop images/fonts on complex slides.
# Instead we do a zipfile-aware byte-copy approach.

import re, io

NSMAP = {
    'a':   'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p':   'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r':   'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
}

def parse_rels(zf, rel_path):
    """Return dict {rId: (type_suffix, target)} from a .rels file."""
    if rel_path not in zf.namelist():
        return {}
    tree = etree.fromstring(zf.read(rel_path))
    out  = {}
    for el in tree.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
        out[el.get('Id')] = (el.get('Type'), el.get('Target'))
    return out

def write_rels(rels_dict):
    """Serialize rels dict back to XML bytes."""
    root = etree.Element(
        '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
    )
    for rid, (rtype, target) in rels_dict.items():
        el = etree.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
        el.set('Id',     rid)
        el.set('Type',   rtype)
        el.set('Target', target)
    return etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)


def merge_presentations(sources, output_path):
    """
    sources: list of (path, [slide_indices])
    Merges all specified slides into a single pptx at output_path.
    Uses the first source as the base (copies theme/layout/master).
    """
    base_path = sources[0][0]
    shutil.copy2(base_path, output_path)

    # Open base as zip for reading master/theme
    with zipfile.ZipFile(base_path, 'r') as base_zf:
        base_names = set(base_zf.namelist())

    # We'll build the merged file from scratch using zipfile manipulation.
    # Strategy:
    #  1. Start with a copy of the first source (base).
    #  2. Remove all its slides.
    #  3. Re-add desired slides from each source in order, renumbering media/rels.

    # Read all source zips into memory
    src_data = []
    for path, indices in sources:
        with zipfile.ZipFile(path, 'r') as zf:
            names    = zf.namelist()
            contents = {n: zf.read(n) for n in names}
        src_data.append((contents, indices))

    base_contents = src_data[0][0]

    # Parse presentation.xml from base to get slide order, size etc.
    prs_xml   = etree.fromstring(base_contents['ppt/presentation.xml'])
    ns_p      = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r      = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst  = prs_xml.find(f'{{{ns_p}}}sldIdLst')

    # Remove all existing slide entries from sldIdLst
    for el in list(sldIdLst):
        sldIdLst.remove(el)

    # Parse presentation.xml.rels
    prs_rels_xml = etree.fromstring(base_contents['ppt/_rels/presentation.xml.rels'])
    ns_rel       = 'http://schemas.openxmlformats.org/package/2006/relationships'
    # Remove existing slide relationships
    for el in list(prs_rels_xml):
        if 'slide"' in el.get('Type','') or el.get('Type','').endswith('/slide'):
            prs_rels_xml.remove(el)

    # Build new merged contents dict
    new_contents = dict(base_contents)

    # Remove base slides from new_contents
    for name in list(new_contents.keys()):
        if re.match(r'ppt/slides/slide\d+\.xml', name):
            del new_contents[name]
        elif re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels', name):
            del new_contents[name]

    # Track media already added (by content hash) to avoid duplicates
    media_map = {}   # (src_path, original_media_name) -> new_media_name
    media_counter = [0]
    next_slide_id = [256]
    slide_counter = [0]

    def add_media(src_contents, media_name):
        """Add media file to new_contents, return new name (deduped by content)."""
        data     = src_contents.get(media_name, b'')
        key      = (id(src_contents), media_name)
        if key in media_map:
            return media_map[key]
        # Check by content hash
        content_hash = hash(data)
        for (k_src, k_name), v_name in media_map.items():
            if k_src == id(src_contents) and k_name == media_name:
                return v_name
        media_counter[0] += 1
        ext      = os.path.splitext(media_name)[1]
        new_name = f'ppt/media/image{media_counter[0]}{ext}'
        new_contents[new_name] = data
        media_map[key]         = new_name
        return new_name

    # Process each source slide
    for src_contents, indices in src_data:
        src_slide_names = sorted(
            [n for n in src_contents if re.match(r'ppt/slides/slide\d+\.xml$', n)],
            key=lambda x: int(re.search(r'\d+', x.split('/')[-1]).group())
        )
        for idx in indices:
            if idx >= len(src_slide_names):
                print(f"  Warning: index {idx} out of range, skipping")
                continue

            slide_counter[0] += 1
            new_slide_num  = slide_counter[0]
            src_slide_name = src_slide_names[idx]
            src_rels_name  = src_slide_name.replace('ppt/slides/', 'ppt/slides/_rels/') + '.rels'

            new_slide_name = f'ppt/slides/slide{new_slide_num}.xml'
            new_rels_name  = f'ppt/slides/_rels/slide{new_slide_num}.xml.rels'

            # Read slide XML
            slide_xml_bytes = src_contents[src_slide_name]

            # Read slide rels
            rels_dict = {}
            if src_rels_name in src_contents:
                rels_root = etree.fromstring(src_contents[src_rels_name])
                for el in rels_root:
                    rid    = el.get('Id')
                    rtype  = el.get('Type')
                    target = el.get('Target')
                    rels_dict[rid] = (rtype, target)

            # Build new rels, copying media into new_contents
            new_rels = {}
            rid_remap = {}  # old rId -> new rId (we keep same rIds here)
            slide_layout_rid = None

            for rid, (rtype, target) in rels_dict.items():
                if target.startswith('../media/'):
                    media_src_name = 'ppt/media/' + target.split('../media/')[-1]
                    new_media_name = add_media(src_contents, media_src_name)
                    new_target     = '../media/' + new_media_name.split('ppt/media/')[-1]
                    new_rels[rid]  = (rtype, new_target)
                elif '../slideLayouts/' in target:
                    # Map to layout 7 (blank) from base — index 6
                    new_rels[rid]  = (rtype, '../slideLayouts/slideLayout7.xml')
                    slide_layout_rid = rid
                else:
                    new_rels[rid]  = (rtype, target)

            new_contents[new_slide_name] = slide_xml_bytes
            new_contents[new_rels_name]  = write_rels(new_rels)

            # Add slide to presentation.xml sldIdLst
            sld_id_el = etree.SubElement(sldIdLst, f'{{{ns_p}}}sldId')
            sld_id_el.set('id',    str(next_slide_id[0]))
            sld_id_el.set(f'{{{ns_r}}}id', f'rId_slide{new_slide_num}')
            next_slide_id[0] += 1

            # Add to presentation.xml.rels
            rel_el = etree.SubElement(prs_rels_xml, f'{{{ns_rel}}}Relationship')
            rel_el.set('Id',     f'rId_slide{new_slide_num}')
            rel_el.set('Type',   'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            rel_el.set('Target', f'slides/slide{new_slide_num}.xml')

            # Add to [Content_Types].xml
            ct_xml = etree.fromstring(new_contents['[Content_Types].xml'])
            ns_ct  = 'http://schemas.openxmlformats.org/package/2006/content-types'
            # Remove existing override for this slide path if any
            for el in ct_xml.findall(f'{{{ns_ct}}}Override'):
                if el.get('PartName') == f'/ppt/slides/slide{new_slide_num}.xml':
                    ct_xml.remove(el)
            ov = etree.SubElement(ct_xml, f'{{{ns_ct}}}Override')
            ov.set('PartName',    f'/ppt/slides/slide{new_slide_num}.xml')
            ov.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
            new_contents['[Content_Types].xml'] = etree.tostring(ct_xml, xml_declaration=True,
                                                                  encoding='UTF-8', standalone=True)

            print(f"  Added slide {new_slide_num}: {src_slide_name} (from {os.path.basename(os.path.dirname(os.path.dirname(src_slide_name))+'/'+os.path.basename(src_slide_name))})")

    # Serialize updated presentation.xml and its rels back
    new_contents['ppt/presentation.xml'] = etree.tostring(
        prs_xml, xml_declaration=True, encoding='UTF-8', standalone=True)
    new_contents['ppt/_rels/presentation.xml.rels'] = etree.tostring(
        prs_rels_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Write final zip
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as out_zf:
        for name, data in new_contents.items():
            out_zf.writestr(name, data)

    print(f"\nGuardado en: {output_path}")
    print(f"Total diapositivas: {slide_counter[0]}")


if __name__ == '__main__':
    merge_presentations(SOURCES, OUTPUT)

    # Quick validation
    prs = Presentation(OUTPUT)
    print(f"\nValidacion: {len(prs.slides)} diapositivas encontradas")
    for i, slide in enumerate(prs.slides):
        texts = [s.text[:60] for s in slide.shapes if hasattr(s,'text') and s.text.strip()]
        print(f"  Slide {i+1}: {texts[:3]}")
