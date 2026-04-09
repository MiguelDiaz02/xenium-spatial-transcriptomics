"""
Presentación Xenium (3 diapositivas) — python-pptx + bioicons
Usa íconos SVG de bioicons.com para enriquecer visualmente cada slide.
"""

import io
import os
import urllib.parse
import requests
import cairosvg
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colores
# ---------------------------------------------------------------------------
BG_NAVY   = RGBColor(0x0A, 0x0F, 0x2C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CYAN      = RGBColor(0x00, 0xD4, 0xFF)
GOLD      = RGBColor(0xFF, 0xD1, 0x66)
GRAY_BLUE = RGBColor(0x8A, 0x9B, 0xBF)
DARK_BOX  = RGBColor(0x14, 0x1E, 0x45)
TEAL      = RGBColor(0x00, 0x7A, 0x8A)

BASE_URL = "https://raw.githubusercontent.com/duerrsimon/bioicons/main/static/icons"

# ---------------------------------------------------------------------------
# Cache local de íconos (evita re-descargar en cada ejecución)
# ---------------------------------------------------------------------------
ICON_CACHE = {}


def load_bioicon(name, category, license_type, author, size_px=200):
    """
    Descarga un SVG de bioicons, lo convierte a PNG en memoria y devuelve BytesIO.
    Cachea los resultados para evitar re-descargar.
    """
    key = (name, size_px)
    if key in ICON_CACHE:
        ICON_CACHE[key].seek(0)
        return ICON_CACHE[key]

    author_enc = urllib.parse.quote(author)
    url = f"{BASE_URL}/{license_type}/{category}/{author_enc}/{name}.svg"
    print(f"  Descargando icono: {name}...")
    r = requests.get(url, timeout=15)
    if r.status_code != 200:
        print(f"  ADVERTENCIA: {name} no encontrado (HTTP {r.status_code}), usando placeholder.")
        return None

    try:
        png_bytes = cairosvg.svg2png(
            bytestring=r.content,
            output_width=size_px,
            output_height=size_px,
            background_color="rgba(0,0,0,0)"  # fondo transparente
        )
        buf = io.BytesIO(png_bytes)
        ICON_CACHE[key] = buf
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"  ERROR convirtiendo {name}: {e}")
        return None


# ---------------------------------------------------------------------------
# Íconos a usar (nombre, categoría, licencia, autor)
# ---------------------------------------------------------------------------
ICONS = {
    "cell":         ("cell-complete",             "Intracellular_components", "cc-by-3.0", "Servier"),
    "tissue":       ("epithelium-squamos",         "Tissues",                 "cc-by-3.0", "Servier"),
    "rna":          ("rna",                        "Nucleic_acids",           "cc-by-3.0", "Servier"),
    "dna":          ("dna-double-stranded-ribbon", "Nucleic_acids",           "cc-by-3.0", "Servier"),
    "antibody":     ("antibody-1",                 "Blood_Immunology",        "cc-by-3.0", "Servier"),
    "microscope":   ("microscope",                 "Lab_apparatus",           "cc-by-3.0", "Servier"),
    "scatter":      ("scatter",                    "Scientific_graphs",       "mit",       "dbs-sticky"),
    "dot":          ("dot-strip",                  "Scientific_graphs",       "mit",       "dbs-sticky"),
    "nucleus":      ("nucleus-closeup",            "Intracellular_components","cc-by-3.0", "Servier"),
    "nn":           ("neural-network-1",           "Machine_Learning",        "cc-0",      "Simon_D\u00fcrr"),
    "classify":     ("classification",             "Machine_Learning",        "cc-0",      "Simon_D\u00fcrr"),
    "liver":        ("liver",                      "Human_physiology",        "cc-by-3.0", "Servier"),
}


def get_icon(key, size_px=200):
    if key not in ICONS:
        return None
    return load_bioicon(*ICONS[key], size_px=size_px)


# ---------------------------------------------------------------------------
# Helpers de diapositiva
# ---------------------------------------------------------------------------

def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=17, color=WHITE, bullet_color=CYAN):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "\u25b8 "
        r1.font.size = Pt(font_size - 1)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
    return txBox


def add_icon_box(slide, icon_key, label, left, top, icon_size=1.2,
                 bg_color=DARK_BOX, border_color=CYAN, label_color=WHITE,
                 label_size=11, icon_px=220):
    """
    Dibuja una caja con ícono centrado arriba + etiqueta de texto abajo.
    """
    box_w = icon_size + 0.15
    box_h = icon_size + 0.55

    # Fondo de la caja
    add_rect(slide, left, top, box_w, box_h,
             fill_rgb=bg_color, line_rgb=border_color, line_width_pt=0.6)

    # Ícono centrado
    img_buf = get_icon(icon_key, size_px=icon_px)
    if img_buf:
        padding = 0.075
        slide.shapes.add_picture(
            img_buf,
            Inches(left + padding),
            Inches(top + 0.06),
            width=Inches(icon_size),
            height=Inches(icon_size)
        )

    # Etiqueta abajo
    add_textbox(slide, label,
                left, top + icon_size + 0.07,
                box_w, 0.42,
                font_size=label_size, bold=True,
                color=label_color, align=PP_ALIGN.CENTER)

    return box_w, box_h


def add_icon_pipeline(slide, steps, left_start, top, icon_size=1.15,
                      box_colors=None, border_color=CYAN, label_color=WHITE,
                      label_size=10, icon_px=200, arrow_color=CYAN):
    """
    Dibuja una fila de cajas de íconos con flechas entre ellas.
    steps: list of (icon_key, label)
    """
    spacing = 0.14
    box_w = icon_size + 0.15

    for i, (icon_key, label) in enumerate(steps):
        x = left_start + i * (box_w + spacing)
        bg = box_colors[i] if box_colors else DARK_BOX
        bd = GOLD if i == len(steps) - 1 else border_color

        add_icon_box(slide, icon_key, label, x, top,
                     icon_size=icon_size,
                     bg_color=bg, border_color=bd,
                     label_color=label_color, label_size=label_size,
                     icon_px=icon_px)

        if i < len(steps) - 1:
            ax = x + box_w
            ay = top + (icon_size + 0.55) / 2 - 0.15
            add_textbox(slide, "\u2192", ax, ay, spacing + 0.02, 0.3,
                        font_size=14, bold=True, color=arrow_color,
                        align=PP_ALIGN.CENTER)

    total_w = len(steps) * box_w + (len(steps) - 1) * spacing
    return total_w


# ---------------------------------------------------------------------------
# Paleta de colores para pipelines
# ---------------------------------------------------------------------------
PIPE_COLORS_TECH = [
    RGBColor(0x1A, 0x2A, 0x55),
    RGBColor(0x00, 0x5F, 0x73),
    RGBColor(0x00, 0x7A, 0x8A),
    RGBColor(0x00, 0x94, 0xA0),
    RGBColor(0x00, 0xB4, 0xC8),
]

PIPE_COLORS_COMP = [
    RGBColor(0x1A, 0x2A, 0x55),
    RGBColor(0x00, 0x55, 0x60),
    RGBColor(0x00, 0x6E, 0x6E),
    RGBColor(0x1A, 0x56, 0x30),
    RGBColor(0x14, 0x45, 0x28),
]


# ---------------------------------------------------------------------------
# Diapositiva 1: El Problema + La Solución
# ---------------------------------------------------------------------------

def build_slide1(slide):
    set_background(slide, BG_NAVY)
    add_rect(slide, 0, 0, 13.33, 0.07, fill_rgb=CYAN)
    add_textbox(slide, "1 / 3", 12.6, 7.1, 0.6, 0.25,
                font_size=11, color=GRAY_BLUE, align=PP_ALIGN.RIGHT)

    add_textbox(slide,
                "Xenium: cada gen (y prote\u00edna) tiene una direcci\u00f3n",
                0.4, 0.15, 9.5, 0.85,
                font_size=30, bold=True, color=CYAN)

    add_textbox(slide,
                "Transcript\u00f3mica espacial de alta resoluci\u00f3n in situ",
                0.4, 0.95, 9.5, 0.4,
                font_size=18, color=GOLD)

    add_rect(slide, 0.4, 1.38, 8.5, 0.025, fill_rgb=GRAY_BLUE)

    bullets = [
        "scRNA-seq: disociamos el tejido \u2192 perdemos la ubicaci\u00f3n celular",
        "La funci\u00f3n de una c\u00e9lula depende de su vecindario \u2014 la posici\u00f3n importa",
        "Xenium: ~5,000 genes + prote\u00ednas detectados in situ, resoluci\u00f3n subcelular",
    ]
    add_bullet_list(slide, bullets, 0.4, 1.55, 7.1, 3.1, font_size=17)

    # --- Visual comparativo con íconos ---
    # Fondo del panel comparativo
    add_rect(slide, 7.7, 1.42, 5.35, 3.35, fill_rgb=RGBColor(0x07, 0x0C, 0x20),
             line_rgb=DARK_BOX, line_width_pt=0.5)

    # Etiqueta superior
    add_textbox(slide, "scRNA-seq  vs  Xenium",
                7.7, 1.45, 5.35, 0.32,
                font_size=12, bold=True, color=GRAY_BLUE, align=PP_ALIGN.CENTER)

    # Caja izquierda: scRNA-seq (célula aislada)
    add_rect(slide, 7.85, 1.82, 2.2, 2.7, fill_rgb=RGBColor(0x25, 0x25, 0x38),
             line_rgb=GRAY_BLUE, line_width_pt=0.4)
    add_textbox(slide, "scRNA-seq", 7.85, 1.87, 2.2, 0.3,
                font_size=11, bold=True, color=GRAY_BLUE, align=PP_ALIGN.CENTER)

    img_cell = get_icon("cell", size_px=200)
    if img_cell:
        slide.shapes.add_picture(img_cell, Inches(8.05), Inches(2.22),
                                 width=Inches(1.8), height=Inches(1.8))
    add_textbox(slide, "Sin contexto\nespacial", 7.85, 4.08, 2.2, 0.38,
                font_size=10, color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

    # Etiqueta VS
    add_textbox(slide, "VS", 10.08, 2.9, 0.35, 0.35,
                font_size=14, bold=True, color=GRAY_BLUE, align=PP_ALIGN.CENTER)

    # Caja derecha: Xenium (tejido estructurado)
    add_rect(slide, 10.45, 1.82, 2.35, 2.7, fill_rgb=RGBColor(0x08, 0x18, 0x35),
             line_rgb=CYAN, line_width_pt=0.8)
    add_textbox(slide, "Xenium", 10.45, 1.87, 2.35, 0.3,
                font_size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    img_tissue = get_icon("tissue", size_px=200)
    if img_tissue:
        slide.shapes.add_picture(img_tissue, Inches(10.65), Inches(2.22),
                                 width=Inches(1.95), height=Inches(1.8))
    add_textbox(slide, "Coordenadas\nprecisas (X, Y)", 10.45, 4.08, 2.35, 0.38,
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Nota técnica y autores
    add_textbox(slide,
                "Formato de datos: SpatialData (Zarr)  \u00b7  Pipeline: recode_st",
                0.4, 4.95, 9.0, 0.3,
                font_size=11, color=GRAY_BLUE, italic=True)

    add_rect(slide, 0.4, 5.35, 12.5, 0.025, fill_rgb=DARK_BOX)
    add_textbox(slide,
                "Miguel \u00c1ngel D\u00edaz Campos  \u00b7  Dr. Alfredo de Jes\u00fas Rodr\u00edguez G\u00f3mez",
                0.4, 5.45, 12.5, 0.35,
                font_size=13, color=GRAY_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "10x Genomics Xenium In Situ  \u00b7  2025",
                0.4, 5.8, 12.5, 0.3,
                font_size=12, color=RGBColor(0x55, 0x66, 0x88), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.4, 13.33, 0.1, fill_rgb=DARK_BOX)


# ---------------------------------------------------------------------------
# Diapositiva 2: Tecnología RNA + Proteína
# ---------------------------------------------------------------------------

def build_slide2(slide):
    set_background(slide, BG_NAVY)
    add_rect(slide, 0, 0, 13.33, 0.07, fill_rgb=CYAN)
    add_textbox(slide, "2 / 3", 12.6, 7.1, 0.6, 0.25,
                font_size=11, color=GRAY_BLUE, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "TECNOLOG\u00cdA", 0.4, 0.1, 3.0, 0.3,
                font_size=11, bold=True, color=CYAN)

    add_textbox(slide,
                "Descifrando el tejido: RNA y prote\u00ednas\nmolécula a mol\u00e9cula",
                0.4, 0.4, 12.5, 1.0,
                font_size=28, bold=True, color=WHITE)

    add_rect(slide, 0.4, 1.38, 6.0, 0.025, fill_rgb=CYAN)

    bullets = [
        "Sondas padlock (RNA) y anticuerpos conjugados (prote\u00ednas)\n    se unen a sus blancos en cortes de tejido intacto",
        "Amplificaci\u00f3n en c\u00edrculo rodante (RCA): las sondas hibridadas\n    forman 'nanobolas' de ADN fluorescentes y detectables",
        "Ciclos de imagen: 4 colores fluorescentes \u00d7 m\u00faltiples rondas\n    \u2192 c\u00f3digo \u00fanico de barras para cada gen / prote\u00edna",
        "Resultado: millones de transcritos y prote\u00ednas con\n    coordenadas (X, Y) precisas en el tejido",
    ]
    add_bullet_list(slide, bullets, 0.4, 1.52, 7.0, 4.1, font_size=16)

    # --- Panel derecho: pipeline de íconos ---
    add_textbox(slide, "Flujo de detecci\u00f3n Xenium",
                7.6, 1.38, 5.3, 0.35,
                font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Fila 1 (3 pasos): Tejido → RNA → Proteína
    steps_row1 = [
        ("tissue",    "Tejido"),
        ("rna",       "RNA (sondas)"),
        ("antibody",  "Prote\u00edna (Ab)"),
    ]
    add_icon_pipeline(slide, steps_row1,
                      left_start=7.65, top=1.82,
                      icon_size=1.3, icon_px=240,
                      box_colors=PIPE_COLORS_TECH[:3],
                      label_size=11)

    # Flecha de bajada
    add_textbox(slide, "\u21d3", 9.85, 3.26, 0.5, 0.3,
                font_size=16, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Fila 2 (2 pasos): Ciclos imagen → Mapa
    steps_row2 = [
        ("microscope", "Ciclos imagen"),
        ("scatter",    "Mapa molecular"),
    ]
    add_icon_pipeline(slide, steps_row2,
                      left_start=8.5, top=3.6,
                      icon_size=1.3, icon_px=240,
                      box_colors=PIPE_COLORS_TECH[3:],
                      border_color=CYAN, label_size=11)

    # Recuadro Xenium Protein
    add_rect(slide, 0.4, 5.3, 12.5, 0.75, fill_rgb=RGBColor(0x1A, 0x14, 0x05),
             line_rgb=GOLD, line_width_pt=1.0)

    # Ícono anticuerpo en el recuadro
    img_ab = get_icon("antibody", size_px=80)
    if img_ab:
        slide.shapes.add_picture(img_ab, Inches(0.55), Inches(5.36),
                                 width=Inches(0.55), height=Inches(0.55))

    add_textbox(slide,
                "\u2605  Xenium Protein (10x Genomics, 2024): anticuerpos conjugados\n"
                "    detectan prote\u00ednas simult\u00e1neamente con el RNA \u2014 mismo tejido, sin procesamiento adicional.",
                1.2, 5.35, 11.5, 0.65,
                font_size=13, color=GOLD)

    add_rect(slide, 0, 7.4, 13.33, 0.1, fill_rgb=DARK_BOX)


# ---------------------------------------------------------------------------
# Diapositiva 3: Del píxel a la biología
# ---------------------------------------------------------------------------

def build_slide3(slide):
    set_background(slide, BG_NAVY)
    add_rect(slide, 0, 0, 13.33, 0.07, fill_rgb=CYAN)
    add_textbox(slide, "3 / 3", 12.6, 7.1, 0.6, 0.25,
                font_size=11, color=GRAY_BLUE, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "AN\u00c1LISIS COMPUTACIONAL", 0.4, 0.1, 5.0, 0.3,
                font_size=11, bold=True, color=CYAN)

    add_textbox(slide, "Del p\u00edxel a la identidad celular",
                0.4, 0.4, 12.5, 0.75,
                font_size=30, bold=True, color=WHITE)

    add_rect(slide, 0.4, 1.15, 6.0, 0.025, fill_rgb=CYAN)

    bullets = [
        "Segmentaci\u00f3n celular (DAPI + membrana) \u2192 transcritos y prote\u00ednas\n    asignados a cada c\u00e9lula individualmente",
        "Matriz de expresi\u00f3n g\u00e9nica \u2192 PCA \u00b7 UMAP \u00b7 Clustering Leiden\n    (reducci\u00f3n de dimensiones e identificaci\u00f3n de grupos celulares)",
        "Anotaci\u00f3n de tipos celulares con marcadores espec\u00edficos\n    de tejido (h\u00edgado, pulm\u00f3n...)",
        "An\u00e1lisis espacial: composici\u00f3n del vecindario celular,\n    gradientes y estad\u00edsticas (Squidpy \u00b7 MuSpAn)",
    ]
    add_bullet_list(slide, bullets, 0.4, 1.3, 7.2, 3.9, font_size=16)

    # --- Pipeline computacional con íconos (derecha) ---
    add_textbox(slide, "Pipeline recode_st (este repositorio)",
                7.55, 1.15, 5.45, 0.35,
                font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Fila 1 (3 pasos): Imagen → Segmentación → Matriz
    steps_row1 = [
        ("microscope", "Imagen\nXenium"),
        ("nucleus",    "Segmen-\ntaci\u00f3n"),
        ("dot",        "Matriz\nRNA+Prot"),
    ]
    add_icon_pipeline(slide, steps_row1,
                      left_start=7.6, top=1.55,
                      icon_size=1.2, icon_px=220,
                      box_colors=PIPE_COLORS_COMP[:3],
                      label_size=10)

    # Flecha bajada
    add_textbox(slide, "\u21d3", 9.8, 2.97, 0.5, 0.3,
                font_size=16, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Fila 2 (2 pasos): UMAP → Mapa espacial
    steps_row2 = [
        ("classify", "UMAP /\nClustering"),
        ("liver",    "Mapa\nEspacial"),
    ]
    add_icon_pipeline(slide, steps_row2,
                      left_start=8.45, top=3.3,
                      icon_size=1.2, icon_px=220,
                      box_colors=PIPE_COLORS_COMP[3:],
                      label_size=10)

    # Módulos del pipeline
    add_rect(slide, 7.55, 4.72, 5.45, 0.55,
             fill_rgb=RGBColor(0x08, 0x0C, 0x20),
             line_rgb=DARK_BOX, line_width_pt=0.5)
    add_textbox(slide,
                "M\u00f3dulos: format_data \u2192 qc \u2192 dimension_reduction\n\u2192 annotate \u2192 spatial_statistics \u00b7 muspan",
                7.65, 4.74, 5.25, 0.5,
                font_size=10, color=GRAY_BLUE, italic=True)

    # Mensaje final
    add_rect(slide, 0.4, 5.42, 12.5, 0.7, fill_rgb=RGBColor(0x05, 0x14, 0x05),
             line_rgb=CYAN, line_width_pt=0.8)
    add_textbox(slide,
                "\u2714  Xenium integra imagen + expresi\u00f3n g\u00e9nica + prote\u00ednas en un solo experimento,\n"
                "    con pipeline reproducible de c\u00f3digo abierto \u2014 listos para usar en su tejido de inter\u00e9s.",
                0.55, 5.47, 12.1, 0.6,
                font_size=13, color=WHITE)

    add_rect(slide, 0, 7.4, 13.33, 0.1, fill_rgb=DARK_BOX)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Construyendo presentacion Xenium con iconos bioicons...")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    print("\n[Diapositiva 1] El Problema + La Solucion")
    s1 = prs.slides.add_slide(blank)
    build_slide1(s1)

    print("\n[Diapositiva 2] Tecnologia RNA + Proteina")
    s2 = prs.slides.add_slide(blank)
    build_slide2(s2)

    print("\n[Diapositiva 3] Del pixel a la biologia")
    s3 = prs.slides.add_slide(blank)
    build_slide3(s3)

    output = "xenium_presentacion.pptx"
    prs.save(output)
    print(f"\nPresentacion guardada: {output}")
    print(f"  Diapositivas: 3 | Dimensiones: 13.33\" x 7.5\" (widescreen)")


if __name__ == "__main__":
    main()
