"""
Presentación Xenium v2 — Diseño inspirado en FlashTalk-MADC
Aplica el sistema de diseño de visual-presentations skill:
  - Banners full-bleed oscuros
  - Columnas con fondo de color completo (Template A)
  - Section strips con pill labels (Template B)
  - Iconos de bioicons como ancla visual dominante (2–2.5")
  - Texto mínimo, jerarquía clara
"""

import io, urllib.parse, requests, cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Dimensiones ──────────────────────────────────────────────────────────────
W = 13.33   # ancho slide (inches)
H = 7.50    # alto slide (inches)

# ─── Paleta FlashTalk + extensión navy ────────────────────────────────────────
NAVY        = RGBColor(0x07, 0x37, 0x63)   # banner primario
TEAL_HDR    = RGBColor(0x13, 0x4F, 0x5C)   # banner secundario
COL_BLUE    = RGBColor(0x11, 0x55, 0xCC)   # columna ciencia
COL_GREEN   = RGBColor(0x1A, 0x6E, 0x3A)   # columna solución (verde oscuro más elegante)
COL_ORANGE  = RGBColor(0xBF, 0x6A, 0x02)   # columna diferenciación (ámbar oscuro)
PILL_BLUE   = RGBColor(0xD9, 0xD2, 0xE9)   # pill azul lavanda
PILL_GREEN  = RGBColor(0xD9, 0xEA, 0xD3)   # pill verde
PILL_ORANGE = RGBColor(0xF9, 0xCB, 0x9C)   # pill naranja/melocotón
SKY_BLUE    = RGBColor(0x6F, 0xA8, 0xDC)   # strip sección A
MAUVE       = RGBColor(0xC2, 0x7B, 0xA0)   # strip sección B
BG_DARK     = RGBColor(0x0A, 0x0F, 0x2C)   # fondo dark opcional
DARK_BOX    = RGBColor(0x14, 0x1E, 0x45)
CYAN        = RGBColor(0x00, 0xD4, 0xFF)
GOLD        = RGBColor(0xFF, 0xD1, 0x66)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY        = RGBColor(0x55, 0x55, 0x66)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
RESULT_PINK = RGBColor(0xF4, 0xCC, 0xCC)

# ─── Bioicons ──────────────────────────────────────────────────────────────────
BASE = "https://raw.githubusercontent.com/duerrsimon/bioicons/main/static/icons"
ICONS = {
    "cell":       ("cell-complete",             "Intracellular_components", "cc-by-3.0", "Servier"),
    "tissue":     ("epithelium-squamos",         "Tissues",                 "cc-by-3.0", "Servier"),
    "rna":        ("rna",                        "Nucleic_acids",           "cc-by-3.0", "Servier"),
    "dna":        ("dna-double-stranded-ribbon", "Nucleic_acids",           "cc-by-3.0", "Servier"),
    "antibody":   ("antibody-1",                 "Blood_Immunology",        "cc-by-3.0", "Servier"),
    "microscope": ("microscope",                 "Lab_apparatus",           "cc-by-3.0", "Servier"),
    "scatter":    ("scatter",                    "Scientific_graphs",       "mit",       "dbs-sticky"),
    "dot":        ("dot-strip",                  "Scientific_graphs",       "mit",       "dbs-sticky"),
    "nucleus":    ("nucleus-closeup",            "Intracellular_components","cc-by-3.0", "Servier"),
    "classify":   ("classification",             "Machine_Learning",        "cc-0",      "Simon_D\u00fcrr"),
    "nn":         ("neural-network-1",           "Machine_Learning",        "cc-0",      "Simon_D\u00fcrr"),
    "liver":      ("liver",                      "Human_physiology",        "cc-by-3.0", "Servier"),
}
_cache = {}

def icon(key, px=300):
    if (key, px) in _cache:
        _cache[(key, px)].seek(0)
        return _cache[(key, px)]
    name, cat, lic, auth = ICONS[key]
    url = f"{BASE}/{lic}/{cat}/{urllib.parse.quote(auth)}/{name}.svg"
    r = requests.get(url, timeout=15)
    if r.status_code != 200:
        print(f"  WARN: {key} → HTTP {r.status_code}")
        return None
    print(f"  Icono: {key} ({px}px)")
    png = cairosvg.svg2png(bytestring=r.content, output_width=px, output_height=px,
                           background_color="rgba(0,0,0,0)")
    buf = io.BytesIO(png)
    _cache[(key, px)] = buf
    return buf

# ─── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def bullets(slide, items, l, t, w, h, size=14, color=WHITE, dot=CYAN):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3); p.space_after = Pt(2)
        r1 = p.add_run(); r1.text = "\u25b8 "; r1.font.size = Pt(size-1)
        r1.font.color.rgb = dot; r1.font.bold = True
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(size); r2.font.color.rgb = color

def pic(slide, key, l, t, w, h=None, px=300):
    buf = icon(key, px=px)
    if buf is None: return
    if h is None: h = w
    slide.shapes.add_picture(buf, Inches(l), Inches(t), width=Inches(w), height=Inches(h))

def banner(slide, title, bg_color=NAVY, text_color=WHITE, size=26):
    """Full-bleed banner ligeramente más ancho que el slide (estilo FlashTalk)."""
    rect(slide, -0.15, 0, W + 0.15, 0.88, fill=bg_color)
    txt(slide, title, 0.25, 0.06, W - 0.5, 0.76,
        size=size, bold=True, color=text_color, align=PP_ALIGN.LEFT)

def section_strip(slide, label, top, color, text_color=WHITE, icon_key=None):
    """Barra de sección estilo FlashTalk (full-width, ~0.40" de alto)."""
    rect(slide, 0, top, W, 0.42, fill=color)
    x_txt = 0.55 if icon_key else 0.25
    txt(slide, label, x_txt, top + 0.04, W - 1.0, 0.34,
        size=13, bold=True, color=text_color)

def pill(slide, label, l, t, w, bg_color, text_color=BLACK, size=11):
    """Pill label estilo FlashTalk (cabecera de columna redondeada)."""
    rect(slide, l, t, w, 0.36, fill=bg_color)
    txt(slide, label, l + 0.08, t + 0.03, w - 0.16, 0.30,
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def slide_number(slide, n, total=3):
    txt(slide, f"{n} / {total}", W - 0.8, H - 0.42, 0.7, 0.32,
        size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def bottom_bar(slide):
    rect(slide, 0, H - 0.12, W, 0.12, fill=DARK_BOX)

# ─── DIAPOSITIVA 1 ─────────────────────────────────────────────────────────────
# Template A: Tres columnas de color completo (patrón FlashTalk)
# Columna 1 Azul  — El problema (scRNA-seq)
# Columna 2 Verde — La solución (Xenium)
# Columna 3 Ámbar — ¿Por qué convencerse?

def slide1(slide):
    bg(slide, WHITE)  # fondo blanco como FlashTalk
    banner(slide, "Xenium: cada gen (y prote\u00edna) tiene una direcci\u00f3n", NAVY)
    slide_number(slide, 1)

    col_y   = 0.92
    col_h   = H - col_y - 0.55   # 6.03"
    col_w   = (W - 0.25) / 3     # ≈4.36"
    gap     = 0.04
    cols    = [
        (COL_BLUE,   PILL_BLUE,   BLACK, "El problema"),
        (COL_GREEN,  PILL_GREEN,  BLACK, "La soluci\u00f3n: Xenium"),
        (COL_ORANGE, PILL_ORANGE, BLACK, "\u00bfPor qu\u00e9 Xenium?"),
    ]

    for i, (col_color, pill_color, pill_txt_color, pill_label) in enumerate(cols):
        x = 0.08 + i * (col_w + gap)

        # Fondo de columna completo
        rect(slide, x, col_y, col_w, col_h, fill=col_color)

        # Pill label
        pill(slide, pill_label, x + 0.08, col_y + 0.06, col_w - 0.16,
             pill_color, pill_txt_color)

    # ── COLUMNA 1 (Problema) ────────────────────────
    x1 = 0.08
    # Ícono célula aislada — GRANDE (ancla visual, ~2.4")
    pic(slide, "cell", x1 + 0.55, col_y + 0.52, 3.2, px=400)
    # Texto debajo del ícono
    txt(slide, "Célula aislada", x1 + 0.08, col_y + 2.95, col_w - 0.16, 0.35,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(slide,
        ["scRNA-seq dissocia el tejido",
         "Perdemos la posici\u00f3n celular",
         "Sin contexto espacial"],
        x1 + 0.15, col_y + 3.38, col_w - 0.30, 1.55,
        size=12, color=LIGHT_GRAY, dot=PILL_BLUE)

    # ── COLUMNA 2 (Solución) ─────────────────────────
    x2 = 0.08 + col_w + gap
    # Ícono tejido estructurado — GRANDE
    pic(slide, "tissue", x2 + 0.55, col_y + 0.52, 3.2, px=400)
    txt(slide, "Tejido intacto", x2 + 0.08, col_y + 2.95, col_w - 0.16, 0.35,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(slide,
        ["~5,000 genes detectados in situ",
         "Coordenadas (X, Y) por transcrito",
         "Resoluci\u00f3n subcelular"],
        x2 + 0.15, col_y + 3.38, col_w - 0.30, 1.55,
        size=12, color=LIGHT_GRAY, dot=PILL_GREEN)

    # ── COLUMNA 3 (Diferenciación) ───────────────────
    x3 = 0.08 + 2 * (col_w + gap)
    # Dos íconos lado a lado: RNA + anticuerpo
    pic(slide, "rna",      x3 + 0.20, col_y + 0.52, 1.8, px=280)
    pic(slide, "antibody", x3 + 2.25, col_y + 0.52, 1.8, px=280)
    txt(slide, "RNA + Prote\u00ednas", x3 + 0.08, col_y + 2.50, col_w - 0.16, 0.35,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Caja de dato destacado
    rect(slide, x3 + 0.15, col_y + 2.92, col_w - 0.30, 0.72,
         fill=RESULT_PINK, line=RGBColor(0xBF,0x6A,0x02), lw=1.0)
    txt(slide, "\u2605 Xenium Protein 2024\nRNA + prote\u00edna simult\u00e1neos",
        x3 + 0.25, col_y + 2.97, col_w - 0.50, 0.62,
        size=11, color=BLACK, bold=False)
    bullets(slide,
        ["Sin procesamiento adicional",
         "Mismo corte de tejido",
         "Anticuerpos conjugados"],
        x3 + 0.15, col_y + 3.72, col_w - 0.30, 1.22,
        size=12, color=LIGHT_GRAY, dot=PILL_ORANGE)

    # Barra de autores
    rect(slide, 0, H - 0.52, W, 0.40, fill=RGBColor(0xE8,0xE8,0xF0))
    txt(slide,
        "Miguel \u00c1ngel D\u00edaz Campos  \u00b7  Dr. Alfredo de Jes\u00fas Rodr\u00edguez G\u00f3mez  \u00b7  10x Genomics Xenium  \u00b7  2025",
        0.3, H - 0.50, W - 0.6, 0.36,
        size=11, color=GRAY, align=PP_ALIGN.CENTER)
    bottom_bar(slide)

# ─── DIAPOSITIVA 2 ─────────────────────────────────────────────────────────────
# Template B: Dos secciones horizontales + figura grande ancla derecha
# Sección A (sky blue): detección molecular — íconos grandes en fila
# Sección B (mauve):    Xenium Protein — caja destacada
# Derecha:              ícono scatter/dot como ancla visual dominante

def slide2(slide):
    bg(slide, WHITE)
    banner(slide, "Descifrando el tejido: RNA y prote\u00ednas mol\u00e9cula a mol\u00e9cula", TEAL_HDR)
    slide_number(slide, 2)

    # ─ Columna derecha: figura ancla dominante ─────────────────────────────────
    right_x   = 8.80
    right_w   = W - right_x - 0.15
    anchor_sz = right_w - 0.10   # ≈4.28"

    rect(slide, right_x - 0.05, 0.92, right_w + 0.05, H - 1.08,
         fill=RGBColor(0xF7, 0xF9, 0xFF), line=RGBColor(0xCC,0xCC,0xDD), lw=0.5)
    txt(slide, "Resultado: mapa molecular",
        right_x, 0.97, right_w, 0.38,
        size=12, bold=True, color=TEAL_HDR, align=PP_ALIGN.CENTER)

    # Ícono scatter grande como "resultado visual"
    pic(slide, "scatter", right_x + 0.10, 1.38, anchor_sz - 0.20, px=480)

    txt(slide,
        "Millones de transcritos\ny prote\u00ednas posicionados\nen el tejido",
        right_x, 1.38 + anchor_sz - 0.05, right_w, 0.90,
        size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # ─ Panel izquierdo ──────────────────────────────────────────────────────────
    left_w = right_x - 0.15

    # ── Sección A: Detección del RNA ─────────────────
    section_strip(slide, "A  \u2014  Detecci\u00f3n del ARN (sondas padlock)", 0.92, SKY_BLUE)

    # Tres tarjetas de herramienta estilo FlashTalk
    cards_a = [
        ("tissue",     "Tejido\nintacto",       "Corte fresco\no FFPE"),
        ("rna",        "Sondas\npadlock",        "Se hibridan\nal ARN diana"),
        ("microscope", "Amplif. +\nImagen",      "RCA \u2192 4 colores\n\u00d7 ciclos"),
    ]
    card_w   = (left_w - 0.30) / 3
    card_h   = 2.10    # aumentado de 1.90 para que la descripción quede dentro
    card_y   = 1.40
    for j, (ic, title, desc) in enumerate(cards_a):
        cx = 0.10 + j * (card_w + 0.05)
        # Fondo tarjeta
        rect(slide, cx, card_y, card_w, card_h,
             fill=RGBColor(0xEB,0xF3,0xFF), line=SKY_BLUE, lw=0.8)
        # Pill header tarjeta
        rect(slide, cx, card_y, card_w, 0.36, fill=SKY_BLUE)
        txt(slide, title, cx + 0.05, card_y + 0.02, card_w - 0.10, 0.32,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Ícono dentro de la tarjeta (reducido a 1.30" para dejar espacio a descripción)
        pic(slide, ic, cx + card_w/2 - 0.65, card_y + 0.40, 1.30, px=280)
        # Descripción — dentro de la tarjeta (después del ícono)
        txt(slide, desc, cx + 0.05, card_y + 1.76, card_w - 0.10, 0.30,
            size=10, color=BLACK, align=PP_ALIGN.CENTER)

        # Flecha conectora
        if j < 2:
            ax = cx + card_w
            txt(slide, "\u2192", ax + 0.01, card_y + card_h/2 - 0.16, 0.08, 0.32,
                size=16, bold=True, color=SKY_BLUE, align=PP_ALIGN.CENTER)

    # ── Sección B: Xenium Protein ─────────────────────
    # strip comienza 0.07" después del final de cards_a (card_y + card_h = 3.50)
    section_strip(slide, "B  \u2014  Xenium Protein: prote\u00ednas espaciales simult\u00e1neas", 3.57, MAUVE)

    # Tres tarjetas sección B
    cards_b = [
        ("antibody",  "Anticuerpos\nconjugados",  "Se unen a\nprote\u00edna diana"),
        ("dna",       "Mismo\nm\u00e9todo FISH",       "Mismos ciclos\nde imagen"),
        ("dot",       "Co-detecci\u00f3n\nRNA+Prot",   "Un solo\nexperimento"),
    ]
    card_y2 = 4.04   # strip termina en 3.57+0.42=3.99; cartas desde 4.04
    for j, (ic, title, desc) in enumerate(cards_b):
        cx = 0.10 + j * (card_w + 0.05)
        rect(slide, cx, card_y2, card_w, card_h,
             fill=RGBColor(0xFA,0xF0,0xF6), line=MAUVE, lw=0.8)
        rect(slide, cx, card_y2, card_w, 0.36, fill=MAUVE)
        txt(slide, title, cx + 0.05, card_y2 + 0.02, card_w - 0.10, 0.32,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        pic(slide, ic, cx + card_w/2 - 0.65, card_y2 + 0.40, 1.30, px=280)
        txt(slide, desc, cx + 0.05, card_y2 + 1.76, card_w - 0.10, 0.30,
            size=10, color=BLACK, align=PP_ALIGN.CENTER)
        if j < 2:
            ax = cx + card_w
            txt(slide, "\u2192", ax + 0.01, card_y2 + card_h/2 - 0.16, 0.08, 0.32,
                size=16, bold=True, color=MAUVE, align=PP_ALIGN.CENTER)

    # Nota fuente
    txt(slide, "Fuente: 10x Genomics Xenium In Situ (2024)",
        0.15, H - 0.50, 8.0, 0.32, size=9, color=GRAY, italic=True)
    bottom_bar(slide)

# ─── DIAPOSITIVA 3 ─────────────────────────────────────────────────────────────
# Template C adaptado: Pipeline computacional como figura dominante
# + Tres tarjetas de resultado abajo

def slide3(slide):
    bg(slide, WHITE)
    banner(slide, "Del p\u00edxel a la identidad celular", NAVY)
    slide_number(slide, 3)

    # ─ Sección pipeline ─────────────────────────────────────────────────────────
    section_strip(slide, "Pipeline computacional  \u2014  recode_st (este repositorio)", 0.92,
                  RGBColor(0x1A, 0x2A, 0x55), text_color=LIGHT_GRAY)

    # 5 íconos grandes en fila — el pipeline es la figura dominante
    pipe_steps = [
        ("microscope", "Imagen\nXenium",     "TIFF +\ncoord.csv"),
        ("nucleus",    "Segmen-\ntaci\u00f3n",       "DAPI +\nmembrana"),
        ("dot",        "Matriz\nRNA+Prot",   "Genes \u00d7\nc\u00e9lulas"),
        ("classify",   "UMAP /\nClustering", "Leiden\ncluster"),
        ("liver",      "Mapa\nEspacial",     "Tipos\ncelulares"),
    ]
    n = len(pipe_steps)
    pipe_w  = W - 0.30
    box_w   = (pipe_w - 0.30) / n   # ≈2.49"
    box_h   = 3.00    # aumentado de 2.60 para incluir sublabel dentro de la caja
    pipe_y  = 1.42    # ligeramente ajustado
    gap     = 0.075

    PIPE_COLORS = [
        RGBColor(0x1A, 0x3A, 0x6E),
        RGBColor(0x08, 0x56, 0x72),
        RGBColor(0x10, 0x6B, 0x6B),
        RGBColor(0x14, 0x62, 0x38),
        RGBColor(0x4A, 0x5E, 0x1A),
    ]

    for i, (ic, label, sublabel) in enumerate(pipe_steps):
        bx = 0.15 + i * (box_w + gap)

        # Fondo caja
        is_last = (i == n - 1)
        border_c = GOLD if is_last else CYAN
        rect(slide, bx, pipe_y, box_w, box_h,
             fill=PIPE_COLORS[i], line=border_c, lw=1.2 if is_last else 0.7)

        # Pill header
        rect(slide, bx, pipe_y, box_w, 0.36, fill=border_c)
        txt(slide, label, bx + 0.05, pipe_y + 0.03, box_w - 0.10, 0.30,
            size=11, bold=True, color=BLACK,
            align=PP_ALIGN.CENTER)

        # Ícono grande (ancla visual) — ~2.0" en caja de 2.49"
        icon_sz = box_w - 0.55
        pic(slide, ic, bx + 0.275, pipe_y + 0.42, icon_sz, px=380)

        # Sublabel dentro de la caja (en la parte inferior)
        txt(slide, sublabel, bx + 0.05, pipe_y + box_h - 0.46, box_w - 0.10, 0.40,
            size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

        # Flecha
        if i < n - 1:
            ax = bx + box_w
            txt(slide, "\u2192", ax + 0.005, pipe_y + box_h/2 - 0.18, gap + 0.06, 0.36,
                size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Etiqueta "Resultado final" bajo el último paso
    last_x = 0.15 + (n-1) * (box_w + gap)
    txt(slide, "\u2605 Resultado final",
        last_x, pipe_y + box_h + 0.04, box_w, 0.26,
        size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ─ Tres tarjetas de resultado ───────────────────────────────────────────────
    res_y  = pipe_y + box_h + 0.38
    res_h  = H - res_y - 0.58
    res_w  = (W - 0.50) / 3
    results = [
        (RGBColor(0x6F,0xA8,0xDC), WHITE,
         "Tipos celulares",
         "Hepatocitos, macr\u00f3fagos, c\u00e9lulas\nendoteliales identificados\nen tejido intacto"),
        (RGBColor(0xC2,0x7B,0xA0), WHITE,
         "Expresi\u00f3n espacial",
         "Mapas de expresi\u00f3n g\u00e9nica\ny prote\u00edca con coordenadas\nprecisas por c\u00e9lula"),
        (RGBColor(0x6A,0xA8,0x4F), WHITE,
         "An\u00e1lisis de vecindario",
         "Composici\u00f3n celular,\ngradientes y estad\u00edsticas\nespaciales (Squidpy \u00b7 MuSpAn)"),
    ]
    for k, (pill_c, pill_tc, rtitle, rdesc) in enumerate(results):
        rx = 0.15 + k * (res_w + 0.05)
        rect(slide, rx, res_y, res_w, res_h,
             fill=RGBColor(0xF7,0xF7,0xFF), line=pill_c, lw=0.8)
        rect(slide, rx, res_y, res_w, 0.34, fill=pill_c)
        txt(slide, rtitle, rx + 0.08, res_y + 0.03, res_w - 0.16, 0.28,
            size=11, bold=True, color=pill_tc, align=PP_ALIGN.CENTER)
        txt(slide, rdesc, rx + 0.10, res_y + 0.40, res_w - 0.20, res_h - 0.48,
            size=11, color=BLACK, align=PP_ALIGN.LEFT)

    # Nota módulos
    txt(slide, "M\u00f3dulos: format_data \u2192 qc \u2192 dimension_reduction \u2192 annotate \u2192 spatial_statistics \u00b7 muspan",
        0.15, H - 0.50, W - 0.30, 0.32,
        size=9, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    bottom_bar(slide)

# ─── Main ──────────────────────────────────────────────────────────────────────

def main():
    print("Presentacion Xenium v2 — FlashTalk design system + bioicons")
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    print("\n[Slide 1] Tres columnas de color (Template A)")
    slide1(prs.slides.add_slide(blank))

    print("\n[Slide 2] Dos secciones + figura ancla (Template B)")
    slide2(prs.slides.add_slide(blank))

    print("\n[Slide 3] Pipeline dominante + tarjetas resultado (Template C)")
    slide3(prs.slides.add_slide(blank))

    out = "xenium_presentacion_v3.pptx"
    prs.save(out)
    print(f"\nGuardado: {out}  ({len(prs.slides)} diapositivas, {W}\" x {H}\")")

if __name__ == "__main__":
    main()
